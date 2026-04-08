from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import argparse
import shutil
from tempfile import NamedTemporaryFile

try:
    import yaml
except ImportError as exc:
    raise RuntimeError(
        "PyYAML is required. Install with: python -m pip install pyyaml"
    ) from exc

from pptx import Presentation

from layout_catalouge import load_layout_catalogue
from layout_resolver import find_layout_by_name
from modality_resolver import resolve_layout
from placeholder_writer import (
    debug_placeholders,
    set_body_bullets,
    set_body_paragraph,
    set_object_text,
    set_picture,
    set_title,
)
from schema_validation import validate_deck_structure
from template_loader import load_presentation_from_template, remove_existing_slides
from text_constraints import validate_text_constraints


SUPPORTED_MODALITIES = {
    "title_slide",
    "index_slide",
    "closing_slide",
    "context_statement",
    "problem_framing",
    "hypothesis_success_criteria",
    "options_considered",
    "chosen_approach",
    "architecture_view",
    "evidence_results",
    "learnings_constraints",
    "implications",
    "next_steps",
    "case_study",
    "strategy",
    "prioritisation",
    "operating_model",
    "section_divider",
    "key_metric",
    "four_pillars",
    "quote_slide",
    "ibm_sign_off",
}


def save_presentation_safely(prs: Presentation, output_path: str | Path) -> None:
    output_path = Path(output_path).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        temp_path = Path(tmp.name)

    try:
        prs.save(temp_path)
        if output_path.exists():
            output_path.unlink()
        shutil.move(str(temp_path), str(output_path))
    except PermissionError as exc:
        if temp_path.exists():
            temp_path.unlink(missing_ok=True)
        raise PermissionError(
            f"Could not write to '{output_path}'. "
            f"The file may already be open in PowerPoint or locked by another process."
        ) from exc


def _write_title_slide(slide, fields: dict) -> None:
    title = str(fields["title"]).strip()
    subtitle = str(fields.get("subtitle", "")).strip()

    if subtitle:
        set_title(slide, f"{title}\n{subtitle}")
    else:
        set_title(slide, title)


def _write_index_slide(slide, fields: dict) -> None:
    set_title(slide, fields["title"])
    set_body_bullets(slide, fields.get("sections", []), idx=12)


def _write_closing_slide(slide, fields: dict) -> None:
    title = str(fields.get("title", "Thank you")).strip()
    set_title(slide, title)

    contact = fields.get("contact")
    if contact:
        try:
            set_body_paragraph(slide, str(contact), idx=12)
        except Exception:
            set_title(slide, f"{title}\n{contact}")


def _write_title_text_slide(slide, fields: dict, body_idx: int = 12) -> None:
    set_title(slide, fields["title"])
    body = fields.get("body", [])
    if isinstance(body, list):
        set_body_bullets(slide, body, idx=body_idx)
    else:
        set_body_paragraph(slide, str(body), idx=body_idx)


def _write_two_column_slide(slide, fields: dict, left_idx: int = 13, right_idx: int = 12) -> None:
    set_title(slide, fields["title"])

    left = fields.get("body_left", [])
    right = fields.get("body_right", [])

    if isinstance(left, list):
        set_body_bullets(slide, left, idx=left_idx)
    else:
        set_body_paragraph(slide, str(left), idx=left_idx)

    if isinstance(right, list):
        set_body_bullets(slide, right, idx=right_idx)
    else:
        set_body_paragraph(slide, str(right), idx=right_idx)


def _write_four_points_slide(slide, fields: dict) -> None:
    set_title(slide, fields["title"])
    points = fields.get("columns", fields.get("points", []))
    point_indices = [12, 13, 14, 15]

    for idx, point in zip(point_indices, points):
        set_body_paragraph(slide, str(point), idx=idx)


def _write_half_image_slide(
    slide,
    fields: dict,
    yaml_base: Path,
    body_idx: int,
    image_idx: int,
) -> None:
    set_title(slide, fields["title"])

    if "body" in fields:
        body = fields.get("body", [])
        if isinstance(body, list):
            set_body_bullets(slide, body, idx=body_idx)
        else:
            set_body_paragraph(slide, str(body), idx=body_idx)
    else:
        bullets = []
        if "lead" in fields and fields["lead"]:
            bullets.append(str(fields["lead"]))
        bullets.extend(fields.get("proof_points", []))
        set_body_bullets(slide, bullets, idx=body_idx)

    image = fields.get("image")
    if image:
        image_path = Path(image)
        if not image_path.is_absolute():
            image_path = yaml_base / image_path
        if image_path.exists():
            set_picture(slide, image_path, idx=image_idx, padding_ratio=0.01)


def _write_insight_boxes_slide(slide, fields: dict) -> None:
    set_title(slide, fields["title"])

    intro = fields.get("intro", [])
    if isinstance(intro, list):
        set_body_bullets(slide, intro, idx=13)
    else:
        set_body_paragraph(slide, str(intro), idx=13)

    boxes = fields.get("boxes", [])
    if len(boxes) > 0:
        set_object_text(slide, boxes[0], idx=17)
    if len(boxes) > 1:
        set_object_text(slide, boxes[1], idx=18)
    if len(boxes) > 2:
        set_object_text(slide, boxes[2], idx=19)


def _write_next_steps_boxes_slide(slide, fields: dict) -> None:
    set_title(slide, fields["title"])
    intro = ["The immediate priority is to move from proof to operational hardening"]
    set_body_bullets(slide, intro, idx=13)

    items = []
    for key in ("body_left", "body_right"):
        value = fields.get(key, [])
        if isinstance(value, list):
            items.extend(value)
        elif value:
            items.append(str(value))

    items = items[:3]

    if len(items) > 0:
        set_object_text(slide, items[0], idx=17)
    if len(items) > 1:
        set_object_text(slide, items[1], idx=18)
    if len(items) > 2:
        set_object_text(slide, items[2], idx=19)


def _write_quote_slide(slide, fields: dict) -> None:
    quote = str(fields.get("quote", "")).strip()
    attribution = str(fields.get("attribution", "")).strip()
    full_text = f"{quote}\n\u2014 {attribution}" if attribution else quote
    set_title(slide, full_text)


def _write_case_study_slide(slide, fields: dict, yaml_base: Path) -> None:
    set_title(slide, fields["title"])

    left = fields.get("body_left", [])
    right = fields.get("body_right", [])

    if isinstance(left, list):
        set_body_bullets(slide, left, idx=13)
    else:
        set_body_paragraph(slide, str(left), idx=13)

    if isinstance(right, list):
        set_body_bullets(slide, right, idx=15)
    else:
        set_body_paragraph(slide, str(right), idx=15)

    image = fields.get("image")
    if image:
        image_path = Path(image)
        if not image_path.is_absolute():
            image_path = yaml_base / image_path
        if image_path.exists():
            set_picture(slide, image_path, idx=14, padding_ratio=0.01)


def add_slide_from_spec(
    prs: Presentation,
    slide_spec: dict,
    yaml_base: Path,
    registry: dict,
) -> None:
    modality = slide_spec["modality"]
    fields = slide_spec.get("fields", {})

    if modality not in SUPPORTED_MODALITIES:
        raise ValueError(
            f"Modality '{modality}' is not supported. "
            f"Supported modalities: {sorted(SUPPORTED_MODALITIES)}"
        )

    resolved = resolve_layout(modality, fields, registry)
    family_name = resolved["family"]
    layout_id = resolved["layout_id"]
    actual_layout_name = resolved["ppt_layout"]

    layout = find_layout_by_name(prs, actual_layout_name)
    slide = prs.slides.add_slide(layout)

    print(f"\nAdded slide using modality: {modality}")
    print(f"Resolved family: {family_name}")
    print(f"Resolved layout id: {layout_id}")
    print(f"Resolved PowerPoint layout: {actual_layout_name}")
    print(f"Available placeholders: {debug_placeholders(slide)}")

    if layout_id in {"title_slide", "cover_image_1", "cover_image_2", "cover_image_7", "cover_image_8"}:
        _write_title_slide(slide, fields)

    elif layout_id == "index_slide":
        _write_index_slide(slide, fields)

    elif layout_id == "thank_you":
        _write_closing_slide(slide, fields)

    elif layout_id in {"title_text", "title_text_split_background"}:
        _write_title_text_slide(slide, fields, body_idx=12)

    elif layout_id in {
        "title_text_two_columns",
        "title_text_two_columns_diff",
        "title_text_two_narrow_columns",
    }:
        _write_two_column_slide(slide, fields, left_idx=13, right_idx=12)

    elif layout_id == "title_text_four_columns":
        _write_four_points_slide(slide, fields)

    elif layout_id in {
        "title_text_half_image",
        "title_image",
        "title_gray_box_over_images",
        "text_gray_box_over_images",
    }:
        _write_half_image_slide(slide, fields, yaml_base, body_idx=13, image_idx=14)

    elif layout_id == "fact_number_half_image":
        _write_half_image_slide(slide, fields, yaml_base, body_idx=12, image_idx=13)

    elif layout_id == "insight_text_boxes":
        if modality == "next_steps":
            _write_next_steps_boxes_slide(slide, fields)
        else:
            _write_insight_boxes_slide(slide, fields)

    elif layout_id == "fact_number":
        if "body" in fields:
            _write_title_text_slide(slide, fields, body_idx=12)
        else:
            bullets = []
            if "lead" in fields and fields["lead"]:
                bullets.append(str(fields["lead"]))
            bullets.extend(fields.get("proof_points", []))
            _write_title_text_slide(slide, {"title": fields["title"], "body": bullets}, body_idx=12)

    elif layout_id in {"case_study_1", "case_study_2"}:
        _write_case_study_slide(slide, fields, yaml_base)

    elif layout_id == "big_text":
        if "quote" in fields:
            _write_quote_slide(slide, fields)
        else:
            set_title(slide, fields["title"])

    elif layout_id == "divider_standard":
        set_title(slide, fields["title"])

    elif layout_id == "divider_with_contents":
        set_title(slide, fields["title"])
        if "sections" in fields:
            set_body_bullets(slide, fields["sections"], idx=13)

    elif layout_id == "quote_layout":
        _write_quote_slide(slide, fields)

    elif layout_id in {"ibm_sign_off_blue80", "ibm_sign_off_blue60", "ibm_sign_off_black"}:
        pass  # purely graphical brand slide — no text placeholders

    else:
        if "title" in fields:
            set_title(slide, fields["title"])


def render_deck(
    template_path: str,
    yaml_path: str,
    output_path: str,
    catalogue_path: str,
) -> None:
    prs, working_pptx = load_presentation_from_template(template_path)
    print(f"Opened working presentation: {working_pptx}")

    remove_existing_slides(prs)

    with open(yaml_path, "r", encoding="utf-8") as f:
        deck_spec = yaml.safe_load(f)

    validate_deck_structure(deck_spec)
    validate_text_constraints(deck_spec)

    registry = load_layout_catalogue(catalogue_path)
    yaml_base = Path(yaml_path).resolve().parent

    for slide_spec in deck_spec["slides"]:
        add_slide_from_spec(prs, slide_spec, yaml_base, registry)

    save_presentation_safely(prs, output_path)
    print(f"\nGenerated deck: {output_path}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True, help="Path to .potx template")
    parser.add_argument("--input", required=True, help="Path to YAML deck spec")
    parser.add_argument(
        "--catalogue",
        default="visual_family_registry.yaml",
        help="Path to visual family registry YAML",
    )
    parser.add_argument(
        "--output",
        default="pocdeck_test_output.pptx",
        help="Output .pptx path",
    )
    args = parser.parse_args()

    render_deck(args.template, args.input, args.output, args.catalogue)


if __name__ == "__main__":
    main()