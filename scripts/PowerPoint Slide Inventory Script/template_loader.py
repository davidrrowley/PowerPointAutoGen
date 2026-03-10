from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile

from pptx import Presentation


TEMPLATE_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)
PRESENTATION_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)


def convert_potx_to_working_pptx(template_path: Path) -> Path:
    """
    Convert a .potx into a working .pptx by patching the Open XML content type.
    """
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    working_file = template_path.with_suffix(".working.pptx")

    with zipfile.ZipFile(template_path, "r") as zin:
        with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            temp_path = Path(tmp.name)

        with zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                if item.filename == "[Content_Types].xml":
                    text = data.decode("utf-8").replace(
                        TEMPLATE_CONTENT_TYPE,
                        PRESENTATION_CONTENT_TYPE,
                    )
                    data = text.encode("utf-8")

                zout.writestr(item, data)

    shutil.move(temp_path, working_file)
    return working_file


def load_presentation_from_template(template_path: str | Path) -> tuple[Presentation, Path]:
    template = Path(template_path)
    working_pptx = convert_potx_to_working_pptx(template)
    prs = Presentation(working_pptx)
    return prs, working_pptx


def remove_existing_slides(prs: Presentation) -> None:
    """
    Remove any starter slides carried through from the template.
    """
    for i in range(len(prs.slides) - 1, -1, -1):
        rel_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[i]