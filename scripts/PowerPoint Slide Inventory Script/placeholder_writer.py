from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Pt

# IBM brand constants
_IBM_BLUE = RGBColor(0x00, 0x43, 0xCE)
_IBM_NEAR_BLACK = RGBColor(0x21, 0x27, 0x2A)
_IBM_FONT = "Arial"
_BODY_FONT_SIZE = Pt(14)
_BODY_SPACE_BEFORE = Pt(4)
_BODY_LINE_SPACING = 1.2


def _all_placeholders(slide):
    return list(slide.placeholders)


def debug_placeholders(slide) -> list[dict]:
    result = []
    for ph in _all_placeholders(slide):
        fmt = ph.placeholder_format
        result.append(
            {
                "idx": fmt.idx,
                "type": str(fmt.type),
                "name": getattr(ph, "name", None),
                "shape_type": str(ph.shape_type),
                "has_text_frame": hasattr(ph, "text_frame"),
                "left": getattr(ph, "left", None),
                "top": getattr(ph, "top", None),
                "width": getattr(ph, "width", None),
                "height": getattr(ph, "height", None),
            }
        )
    return result


def get_placeholder(slide, ph_type=None, idx: int | None = None):
    for ph in _all_placeholders(slide):
        fmt = ph.placeholder_format

        if ph_type is not None and fmt.type != ph_type:
            continue
        if idx is not None and fmt.idx != idx:
            continue

        return ph

    raise ValueError(
        f"Placeholder not found. Requested type={ph_type}, idx={idx}. "
        f"Available placeholders={debug_placeholders(slide)}"
    )


def get_first_text_placeholder(slide, exclude_title: bool = True):
    candidates = []

    for ph in _all_placeholders(slide):
        fmt = ph.placeholder_format
        ph_type = fmt.type

        if exclude_title and ph_type == PP_PLACEHOLDER.TITLE:
            continue

        if hasattr(ph, "text_frame"):
            candidates.append(ph)

    for ph in candidates:
        if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return ph

    if candidates:
        return candidates[0]

    raise ValueError(
        f"No text-capable placeholder found. Available placeholders={debug_placeholders(slide)}"
    )


def set_title(slide, text: str) -> None:
    ph = get_placeholder(slide, PP_PLACEHOLDER.TITLE)
    ph.text = text


def _force_bullet(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()

    for child in list(pPr):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in {"buNone", "buChar", "buAutoNum", "buBlip"}:
            pPr.remove(child)

    bu = paragraph._element.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
    )
    bu.set("char", "•")
    pPr.insert(0, bu)


def _style_body_paragraph(paragraph) -> None:
    """Apply IBM body typography to a paragraph (spacing + font, no colour override)."""
    paragraph.space_before = _BODY_SPACE_BEFORE
    for run in paragraph.runs:
        run.font.name = _IBM_FONT


def set_body_bullets(slide, bullets, idx: int | None = None) -> None:
    ph = get_placeholder(slide, PP_PLACEHOLDER.BODY, idx=idx)
    tf = ph.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        _force_bullet(p)
        _style_body_paragraph(p)


def set_body_paragraph(slide, text: str, idx: int | None = None) -> None:
    ph = get_placeholder(slide, PP_PLACEHOLDER.BODY, idx=idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    _style_body_paragraph(p)


def set_first_text(slide, text_or_bullets) -> None:
    ph = get_first_text_placeholder(slide, exclude_title=True)
    tf = ph.text_frame
    tf.clear()

    if isinstance(text_or_bullets, list):
        for i, bullet in enumerate(text_or_bullets):
            # Guard: single-key dicts (broken YAML bullets) → "key: value" string
            if isinstance(bullet, dict) and len(bullet) == 1:
                k, v = next(iter(bullet.items()))
                bullet = f"{k}: {v}"
            elif not isinstance(bullet, str):
                bullet = str(bullet)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            _force_bullet(p)
    else:
        tf.paragraphs[0].text = str(text_or_bullets)


def set_object_text(slide, text: str, idx: int) -> None:
    ph = get_placeholder(slide, PP_PLACEHOLDER.OBJECT, idx=idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    _style_body_paragraph(p)


def set_speaker_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = text.strip()


def set_picture(
    slide,
    image_path: str | Path,
    idx: int | None = None,
    padding_ratio: float = 0.01,
) -> None:
    image_path = Path(image_path)
    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    ph = get_placeholder(slide, PP_PLACEHOLDER.PICTURE, idx=idx)

    left = ph.left
    top = ph.top
    box_w = ph.width
    box_h = ph.height

    pad_w = int(box_w * padding_ratio)
    pad_h = int(box_h * padding_ratio)

    inner_left = left + pad_w
    inner_top = top + pad_h
    inner_w = box_w - (2 * pad_w)
    inner_h = box_h - (2 * pad_h)

    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        box_w,
        box_h,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    with Image.open(image_path) as img:
        img_w_px, img_h_px = img.size

    img_ratio = img_w_px / img_h_px
    box_ratio = inner_w / inner_h

    if img_ratio > box_ratio:
        new_w = inner_w
        new_h = int(inner_w / img_ratio)
    else:
        new_h = inner_h
        new_w = int(inner_h * img_ratio)

    new_left = inner_left + int((inner_w - new_w) / 2)
    new_top = inner_top + int((inner_h - new_h) / 2)

    slide.shapes.add_picture(
        str(image_path),
        new_left,
        new_top,
        width=new_w,
        height=new_h,
    )

    sp = ph._element
    sp.getparent().remove(sp)