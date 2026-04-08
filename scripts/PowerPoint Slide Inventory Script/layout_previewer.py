"""
layout_previewer.py
-------------------
Generate a .pptx file with one sample slide per IBM template layout,
labelled with its layout name and placeholder indices.

Useful for visually browsing the full template catalogue and
verifying which layouts work best for each modality.

Usage:
    python layout_previewer.py --template IBM_Consulting_Presentation_Template_2022_V02_Arial.potx
                                --output layout_catalogue_preview.pptx
"""
from __future__ import annotations

import argparse
import shutil
from pathlib import Path
from tempfile import NamedTemporaryFile

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

from template_loader import load_presentation_from_template, remove_existing_slides


_SAMPLE_TITLE = "Layout: {name}"
_SAMPLE_BODY = "idx={idx}: {ph_name}"

_SKIPPED_LAYOUT_KEYWORDS = {
    "Animated", "Static Rebus", "ibm sign-off",
}


def _should_skip(layout_name: str) -> bool:
    return any(kw in layout_name for kw in _SKIPPED_LAYOUT_KEYWORDS)


def _write_sample_slide(prs: Presentation, layout) -> None:
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        fmt = ph.placeholder_format
        ph_type = fmt.type
        idx = fmt.idx

        if not hasattr(ph, "text_frame"):
            continue

        tf = ph.text_frame
        tf.clear()

        if ph_type == PP_PLACEHOLDER.TITLE:
            tf.paragraphs[0].text = _SAMPLE_TITLE.format(name=layout.name)
            tf.paragraphs[0].runs[0].font.size = Pt(20) if len(layout.name) > 30 else Pt(24)

        elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            tf.paragraphs[0].text = _SAMPLE_BODY.format(
                idx=idx, ph_name=getattr(ph, "name", "?")
            )
            if len(tf.paragraphs[0].runs) > 0:
                tf.paragraphs[0].runs[0].font.size = Pt(14)

        else:
            # Footer / slide number placeholders — leave as-is
            pass


def build_catalogue_preview(template_path: str, output_path: str) -> int:
    prs, _ = load_presentation_from_template(template_path)
    remove_existing_slides(prs)

    added = 0
    for layout in prs.slide_layouts:
        if _should_skip(layout.name):
            continue

        ph_info = [
            f"idx={ph.placeholder_format.idx}({str(ph.placeholder_format.type).split('(')[1].rstrip(')')})"
            for ph in layout.placeholders
        ]
        print(f"  {layout.name:55} placeholders: {ph_info}")

        _write_sample_slide(prs, layout)
        added += 1

    output_path = Path(output_path).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        temp_path = Path(tmp.name)

    prs.save(temp_path)
    if output_path.exists():
        output_path.unlink()
    shutil.move(str(temp_path), str(output_path))

    return added


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate a .pptx file showing one sample slide per IBM template layout"
    )
    parser.add_argument(
        "--template",
        default="IBM_Consulting_Presentation_Template_2022_V02_Arial.potx",
        help="Path to .potx template file",
    )
    parser.add_argument(
        "--output",
        default="layout_catalogue_preview.pptx",
        help="Output .pptx path (default: layout_catalogue_preview.pptx)",
    )
    args = parser.parse_args()

    print(f"Building layout catalogue preview from: {args.template}")
    count = build_catalogue_preview(args.template, args.output)
    print(f"\nGenerated {count} sample slides → {Path(args.output).resolve()}")


if __name__ == "__main__":
    main()
