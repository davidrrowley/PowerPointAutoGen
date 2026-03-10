import argparse
import shutil
from pathlib import Path

from pptx import Presentation


import shutil
import zipfile
from pathlib import Path
from tempfile import NamedTemporaryFile


def copy_template_to_pptx(template_path: Path) -> Path:
    """
    Convert a .potx into a real working .pptx by patching the package
    content type from template to presentation.
    """
    working_file = template_path.with_suffix(".working.pptx")

    with zipfile.ZipFile(template_path, "r") as zin:
        with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            temp_path = Path(tmp.name)

        with zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)

                if item.filename == "[Content_Types].xml":
                    text = data.decode("utf-8").replace(
                        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
                    )
                    data = text.encode("utf-8")

                zout.writestr(item, data)

    shutil.move(temp_path, working_file)
    return working_file


def find_layout(prs: Presentation, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name.lower() == layout_name.lower():
            return layout
    raise ValueError(f"Layout '{layout_name}' not found in template.")


def get_placeholder(slide, ph_type=None, idx=None):
    for ph in slide.placeholders:
        if ph_type and ph.placeholder_format.type != ph_type:
            continue
        if idx is not None and ph.placeholder_format.idx != idx:
            continue
        return ph
    raise ValueError("Placeholder not found.")


def add_big_text_slide(prs):
    layout = find_layout(prs, "big text")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:  # TITLE
            ph.text = "Structured PoC slide generation can preserve narrative clarity"

    return slide


def add_problem_slide(prs):
    layout = find_layout(prs, "title, text")
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:  # TITLE
            ph.text = "Manual PoC slide creation introduces inconsistency"

        if ph.placeholder_format.type == 2:  # BODY
            tf = ph.text_frame
            tf.clear()

            bullets = [
                "Narrative structure varies between engagements",
                "Slides often become overly text dense",
                "Formatting effort distracts from reasoning",
                "A repeatable structure would improve clarity",
            ]

            for i, text in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = 0

    return slide


def add_hypothesis_slide(prs):
    layout = find_layout(prs, "title, text (two columns)")
    slide = prs.slides.add_slide(layout)

    title_set = False
    bodies = []

    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "The PoC must prove both narrative discipline and technical feasibility"
            title_set = True

        if ph.placeholder_format.type == 2:
            bodies.append(ph)

    if len(bodies) >= 2:
        left = bodies[0].text_frame
        left.clear()

        left_points = [
            "Can a PoC narrative arc be expressed as structured data?",
            "Can titles and bullets be constrained automatically?",
        ]

        for i, text in enumerate(left_points):
            p = left.paragraphs[0] if i == 0 else left.add_paragraph()
            p.text = text
            p.level = 0

        right = bodies[1].text_frame
        right.clear()

        right_points = [
            "Can PowerPoint templates be inspected reliably?",
            "Can layouts be populated deterministically from Python?",
        ]

        for i, text in enumerate(right_points):
            p = right.paragraphs[0] if i == 0 else right.add_paragraph()
            p.text = text
            p.level = 0

    return slide


def add_architecture_slide(prs, image_path):
    layout = find_layout(prs, "title, text, half-image")
    slide = prs.slides.add_slide(layout)

    body_placeholder = None
    pic_placeholder = None

    for ph in slide.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Architecture: narrative-driven slide generation pipeline"

        if ph.placeholder_format.type == 2:
            body_placeholder = ph

        if ph.placeholder_format.type == 18:  # picture
            pic_placeholder = ph

    if body_placeholder:
        tf = body_placeholder.text_frame
        tf.clear()

        bullets = [
            "Narrative arc expressed as structured input",
            "AI generates constrained slide content",
            "Template layouts enforce spatial discipline",
        ]

        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = 0

    if image_path and pic_placeholder:
        pic_placeholder.insert_picture(str(image_path))

    return slide


def remove_existing_slides(prs):
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]


def main(template, output, image):
    template = Path(template)

    working_pptx = copy_template_to_pptx(template)
    prs = Presentation(working_pptx)

    remove_existing_slides(prs)

    add_big_text_slide(prs)
    add_problem_slide(prs)
    add_hypothesis_slide(prs)

    if image:
        add_architecture_slide(prs, Path(image))

    prs.save(output)

    print(f"\nGenerated test deck: {output}\n")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True)
    parser.add_argument("--output", default="pocdeck_test_output.pptx")
    parser.add_argument("--image", default=None)

    args = parser.parse_args()

    main(args.template, args.output, args.image)